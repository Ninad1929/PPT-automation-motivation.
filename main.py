import os
import shutil
from PIL import Image
import torch
import clip
from pptx import Presentation
from pptx.util import Inches
import re

# -------- CONFIG --------
IMAGE_FOLDER = "motivational_images"    # put all your raw images here
OUTPUT_FOLDER = "Final_images"          # categorized folders will be created here
PPT_OUTPUT_FOLDER = "presentations"     # folder for generated PPTs

CATEGORIES = [
    "Success & Hard Work",
    "Habits & Discipline",
    "Time & Productivity",
    "Attitude & Positivity",
    "Knowledge & Learning",
    "Life Lessons / General Motivation"
]

# PowerPoint slide size
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

# -------- LOAD CLIP MODEL --------
device = "cuda" if torch.cuda.is_available() else "cpu"
model, preprocess = clip.load("ViT-B/32", device=device)

# Encode category labels
text_tokens = clip.tokenize(CATEGORIES).to(device)
with torch.no_grad():
    text_features = model.encode_text(text_tokens)
    text_features /= text_features.norm(dim=-1, keepdim=True)

# -------- STEP 1: AUTO-CATEGORIZE IMAGES --------
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

def safe_folder_name(name: str) -> str:
    """Convert category name into safe Windows folder name"""
    return re.sub(r'[\\/*?:"<>|]', "_", name)

for img_file in os.listdir(IMAGE_FOLDER):
    if img_file.lower().endswith((".jpg", ".jpeg", ".png")):
        img_path = os.path.join(IMAGE_FOLDER, img_file)

        try:
            image = preprocess(Image.open(img_path)).unsqueeze(0).to(device)
        except Exception:
            print(f"⚠️ Skipping {img_file} (cannot open image)")
            continue

        with torch.no_grad():
            image_features = model.encode_image(image)
            image_features /= image_features.norm(dim=-1, keepdim=True)

        similarity = (image_features @ text_features.T).squeeze(0)
        best_idx = similarity.argmax().item()
        best_category = CATEGORIES[best_idx]

        safe_cat = safe_folder_name(best_category)
        category_path = os.path.join(OUTPUT_FOLDER, safe_cat)
        os.makedirs(category_path, exist_ok=True)

        shutil.copy(img_path, os.path.join(category_path, img_file))
        print(f"✅ {img_file} → {best_category} ({safe_cat})")

print("\n🎉 Step 1 done: Images categorized into motivational folders!\n")

# -------- STEP 2: BUILD PPTs FOR EACH CATEGORY --------
os.makedirs(PPT_OUTPUT_FOLDER, exist_ok=True)

for category in CATEGORIES:
    safe_cat = safe_folder_name(category)
    cat_path = os.path.join(OUTPUT_FOLDER, safe_cat)

    if not os.path.isdir(cat_path):
        continue

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]

    # Title slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = category
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = ""

    # Add all images
    for img_file in sorted(os.listdir(cat_path)):
        if img_file.lower().endswith((".jpg", ".jpeg", ".png")):
            slide = prs.slides.add_slide(blank_slide_layout)
            img_path = os.path.join(cat_path, img_file)

            im = Image.open(img_path)
            width, height = im.size
            aspect = width / height

            max_width = Inches(9)
            max_height = Inches(6.5)

            if aspect > 1:  # landscape
                pic_width = max_width
                pic_height = pic_width / aspect
                if pic_height > max_height:
                    pic_height = max_height
                    pic_width = pic_height * aspect
            else:  # portrait
                pic_height = max_height
                pic_width = pic_height * aspect
                if pic_width > max_width:
                    pic_width = max_width
                    pic_height = pic_width / aspect

            left = (SLIDE_WIDTH - pic_width) / 2
            top = (SLIDE_HEIGHT - pic_height) / 2

            slide.shapes.add_picture(img_path, left, top, width=pic_width, height=pic_height)

    ppt_path = os.path.join(PPT_OUTPUT_FOLDER, f"{safe_cat}.pptx")
    prs.save(ppt_path)
    print(f"📊 Saved PPT: {ppt_path}")

print("\n✅ All PPTs generated successfully in 'presentations/' folder!")
